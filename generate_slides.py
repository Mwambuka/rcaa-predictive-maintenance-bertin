"""
Generates the PowerPoint Presentation for RCAA Predictive Maintenance Assessment.
12 slides designed for a 10–15 min talk to mixed technical/non-technical audience.

Story arc:
  1. Title
  2. The Problem (Stakes)
  3. Our Approach (Roadmap)
  4. The Data
  5. What the Data Tells Us (EDA)
  6. Turning Snapshots into Signals (Feature Engineering)
  7. Training Four Models
  8. Results — Zero Missed Failures
  9. Why Does the Model Predict Failure? (Explainability)
 10. Choosing the Right Alert Threshold
 11. From Predictions to Action (Decision Support)
 12. Key Takeaways & Next Steps
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os, json

# ── Colours ───────────────────────────────────────────────────────────────────
DARK_BLUE   = RGBColor(0x1f, 0x4e, 0x79)
MED_BLUE    = RGBColor(0x29, 0x80, 0xb9)
LIGHT_BLUE  = RGBColor(0xd6, 0xe4, 0xf0)
RED         = RGBColor(0xc0, 0x39, 0x2b)
GREEN       = RGBColor(0x27, 0xae, 0x60)
ORANGE      = RGBColor(0xe6, 0x7e, 0x22)
GOLD        = RGBColor(0xf1, 0xc4, 0x0f)
WHITE       = RGBColor(0xff, 0xff, 0xff)
DARK_GRAY   = RGBColor(0x2c, 0x3e, 0x50)
LIGHT_GRAY  = RGBColor(0xf5, 0xf6, 0xfa)

FIGURES  = "figures"
with open("final_results_summary.json") as f:
    res_data = json.load(f)
best_name = res_data["best_model"]

# ── Helpers ───────────────────────────────────────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def rect(slide, l, t, w, h, fill=None, line=None, line_w=Pt(0)):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line; s.line.width = line_w
    else:
        s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h, size=14, bold=False, color=DARK_GRAY,
        align=PP_ALIGN.LEFT, italic=False, font="Calibri"):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf  = box.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text            = text
    run.font.size       = Pt(size)
    run.font.bold       = bold
    run.font.italic     = italic
    run.font.color.rgb  = color
    run.font.name       = font
    return box

def img(slide, path, l, t, w=None, h=None):
    if not os.path.exists(path): return
    if w and h:
        slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))
    elif w:
        slide.shapes.add_picture(path, Inches(l), Inches(t), width=Inches(w))
    else:
        slide.shapes.add_picture(path, Inches(l), Inches(t))

def header(slide, title, subtitle=None):
    rect(slide, 0, 0, 13.33, 1.35, fill=DARK_BLUE)
    rect(slide, 0, 1.35, 13.33, 0.07, fill=MED_BLUE)
    txt(slide, title, 0.3, 0.1, 12.5, 0.75,
        size=28, bold=True, color=WHITE, font="Calibri Light")
    if subtitle:
        txt(slide, subtitle, 0.3, 0.85, 12.5, 0.42,
            size=12, color=LIGHT_BLUE, font="Calibri Light")

def footer_bar(slide, page, total=12):
    rect(slide, 0, 7.15, 13.33, 0.35, fill=DARK_BLUE)
    txt(slide, "RCAA — Predictive Maintenance Assessment  |  March 2026",
        0.2, 7.15, 10.5, 0.35, size=8, color=LIGHT_BLUE)
    txt(slide, f"{page} / {total}", 12.5, 7.15, 0.8, 0.35,
        size=8, color=LIGHT_BLUE, align=PP_ALIGN.RIGHT)

def kpi_box(slide, l, t, value, label, bg=MED_BLUE):
    rect(slide, l, t, 2.9, 1.45, fill=bg)
    txt(slide, value, l, t + 0.1, 2.9, 0.85,
        size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(slide, label, l, t + 0.95, 2.9, 0.4,
        size=10, color=WHITE, align=PP_ALIGN.CENTER)

def bullets(slide, items, l, t, w, h, size=13.5, gap=0.48):
    """items: list of (indent_level, text)"""
    for i, (ind, text) in enumerate(items):
        prefix = "   " * ind + ("▸ " if ind == 0 else "  – ")
        txt(slide, prefix + text, l, t + i * gap, w, gap + 0.05,
            size=size, color=DARK_GRAY)

# ── Build ─────────────────────────────────────────────────────────────────────
prs = new_prs()

# ─────────────────────────────────────────────────────────────────────────────
# 1. TITLE
# ─────────────────────────────────────────────────────────────────────────────
slide = blank(prs)
rect(slide, 0, 0, 13.33, 7.5, fill=DARK_BLUE)
rect(slide, 0, 4.6, 13.33, 2.9, fill=MED_BLUE)
rect(slide, 0, 4.5, 13.33, 0.12, fill=WHITE)
txt(slide, "Predictive Maintenance",
    0.5, 0.8, 12.3, 1.1, size=46, bold=True, color=WHITE,
    align=PP_ALIGN.CENTER, font="Calibri Light")
txt(slide, "for Aircraft Components",
    0.5, 1.75, 12.3, 0.85, size=36, color=LIGHT_BLUE,
    align=PP_ALIGN.CENTER, font="Calibri Light")
txt(slide, "Rwanda Civil Aviation Authority — Data Science Assessment",
    0.5, 2.8, 12.3, 0.5, size=15, color=WHITE, align=PP_ALIGN.CENTER)
txt(slide, "March 2026",
    0.5, 5.05, 12.3, 0.5, size=15, color=WHITE, align=PP_ALIGN.CENTER)
txt(slide, "Predicting component failure within the next 10 flight cycles",
    0.5, 5.7, 12.3, 0.5, size=13, color=LIGHT_BLUE,
    align=PP_ALIGN.CENTER, italic=True)

# ─────────────────────────────────────────────────────────────────────────────
# 2. THE PROBLEM (STAKES)
# ─────────────────────────────────────────────────────────────────────────────
slide = blank(prs)
header(slide, "The Problem We're Solving",
       "Every unplanned component failure is a failure of information")

# Impact boxes
impact_data = [
    (RED,    "AOG Event",        "Aircraft-on-Ground\ncosts $10k–$150k/hr"),
    (ORANGE, "Safety Risk",      "Unplanned failures\ncompromise safety margins"),
    (GOLD,   "Delayed Flights",  "Cascading disruption\nacross the schedule"),
    (MED_BLUE,"Wasted Resource", "Early maintenance\ncosts with no benefit"),
]
for i, (color, title, desc) in enumerate(impact_data):
    rect(slide, 0.25 + i * 3.2, 1.6, 3.0, 2.1, fill=color)
    txt(slide, title, 0.35 + i*3.2, 1.7, 2.8, 0.55,
        size=14, bold=True, color=WHITE)
    txt(slide, desc, 0.35 + i*3.2, 2.3, 2.8, 1.1, size=11, color=WHITE)

rect(slide, 0.25, 3.85, 12.85, 1.45, fill=LIGHT_BLUE)
txt(slide, "The Solution: Predict failure 10 flight cycles in advance",
    0.45, 3.92, 12.3, 0.5, size=14, bold=True, color=DARK_BLUE)
txt(slide, "Give maintenance planners a ranked list of at-risk components every flight — "
    "before any failure occurs, while there is still time to act.",
    0.45, 4.45, 12.3, 0.75, size=12.5, color=DARK_GRAY)

txt(slide, "\"Catching zero failures is not a data science achievement. "
    "Catching every failure before it grounds a plane is.\"",
    0.3, 5.5, 12.7, 0.55, size=11, color=DARK_BLUE,
    align=PP_ALIGN.CENTER, italic=True)
footer_bar(slide, 2)

# ─────────────────────────────────────────────────────────────────────────────
# 3. OUR APPROACH
# ─────────────────────────────────────────────────────────────────────────────
slide = blank(prs)
header(slide, "Our Approach", "Five steps from raw sensor data to actionable risk scores")

steps = [
    (MED_BLUE,  "1",  "Data Quality",    "Validate sensors,\nflag anomalies"),
    (GREEN,     "2",  "Explore (EDA)",   "Understand failure\npatterns and signals"),
    (ORANGE,    "3",  "Engineer Features","Transform snapshots\ninto temporal signals"),
    (RED,       "4",  "Train & Compare", "4 models, cost-sensitive,\ncross-validated"),
    (RGBColor(0x8e,0x44,0xad), "5", "Deploy",
     "Risk scores →\nmaintenance action"),
]
for i, (color, num, title, desc) in enumerate(steps):
    lx = 0.25 + i * 2.6
    rect(slide, lx, 1.6, 2.4, 3.8, fill=color)
    txt(slide, num, lx, 1.65, 2.4, 0.75,
        size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri Light")
    txt(slide, title, lx + 0.1, 2.4, 2.2, 0.6,
        size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(slide, desc, lx + 0.1, 3.1, 2.2, 1.8,
        size=11, color=WHITE, align=PP_ALIGN.CENTER)

# Connecting arrows
for i in range(4):
    txt(slide, "→", 2.48 + i*2.6, 2.85, 0.3, 0.5,
        size=18, bold=True, color=DARK_GRAY, align=PP_ALIGN.CENTER)

txt(slide, "Throughout: explainability (SHAP + decision rules) so engineers understand every alert",
    0.25, 5.6, 12.85, 0.55, size=12, bold=True, color=DARK_BLUE,
    align=PP_ALIGN.CENTER)
footer_bar(slide, 3)

# ─────────────────────────────────────────────────────────────────────────────
# 4. THE DATA
# ─────────────────────────────────────────────────────────────────────────────
slide = blank(prs)
header(slide, "The Dataset", "6,000 records | 20 aircraft | 4 components | 1 critical challenge")

stats = [
    ("6,000",  "Records",          MED_BLUE),
    ("20",     "Aircraft",         MED_BLUE),
    ("4",      "Component Types",  MED_BLUE),
    ("140",    "Failure Events",   RED),
    ("42:1",   "Class Imbalance",  ORANGE),
    ("13",     "Raw Features",     GREEN),
]
for i, (val, lbl, bg) in enumerate(stats):
    row, col = divmod(i, 3)
    kpi_box(slide, 0.25 + col*3.05, 1.55 + row*1.65, val, lbl, bg=bg)

rect(slide, 9.6, 1.55, 3.45, 3.25, fill=LIGHT_BLUE)
txt(slide, "⚠ The Core Challenge", 9.75, 1.62, 3.2, 0.45,
    size=11, bold=True, color=RED)
txt(slide, "Only 2.3% of records\nare failures.\n\n"
    "A model that always\npredicts 'No Failure'\nscores 97.7% accuracy\n—\n"
    "and catches ZERO\nfailures.",
    9.75, 2.1, 3.2, 2.6, size=11, color=DARK_GRAY)

img(slide, f"{FIGURES}/01_target_distribution.png", 0.15, 5.0, w=9.3, h=2.4)
footer_bar(slide, 4)

# ─────────────────────────────────────────────────────────────────────────────
# 5. WHAT THE DATA TELLS US
# ─────────────────────────────────────────────────────────────────────────────
slide = blank(prs)
header(slide, "What the Data Tells Us",
       "No single sensor predicts failure — it's the pattern over time that matters")

img(slide, f"{FIGURES}/05_boxplots_by_failure.png", 0.15, 1.5, w=8.5, h=4.5)

rect(slide, 8.9, 1.5, 4.2, 4.5, fill=LIGHT_BLUE)
txt(slide, "Key Findings", 9.05, 1.6, 3.9, 0.45,
    size=13, bold=True, color=DARK_BLUE)
key_findings = [
    (0, "Fault codes: clear separation"),
    (0, "Vibration elevated pre-failure"),
    (0, "Temperature: mild signal"),
    (0, "Humidity/ambient: no signal"),
    (0, "Engine1 + Wing: highest risk"),
    (0, "Failures cluster at high cycles"),
    (0, "Sensor drift flags real issues"),
]
bullets(slide, key_findings, 8.95, 2.15, 4.05, 3.7, size=12, gap=0.47)

txt(slide, "→ Single-cycle readings hint at the signal. "
    "Tracking accumulation over 5 cycles unlocks it.",
    0.15, 6.1, 12.9, 0.45, size=11.5, bold=True, color=DARK_BLUE,
    align=PP_ALIGN.CENTER, italic=True)
footer_bar(slide, 5)

# ─────────────────────────────────────────────────────────────────────────────
# 6. FEATURE ENGINEERING
# ─────────────────────────────────────────────────────────────────────────────
slide = blank(prs)
header(slide, "Turning Snapshots into Signals",
       "Feature engineering: from 13 raw features to 43 temporal features")

fe_blocks = [
    (MED_BLUE,  "Rolling Stats",      "Mean & std dev\nover 3 & 5 cycles"),
    (GREEN,     "Rate of Change",     "Δ per cycle\n— speed of deterioration"),
    (ORANGE,    "Fault Accumulation", "★ Rolling 5-cycle\nfault total"),
    (RED,       "Stress Index",       "Vibration + Pressure\ncombined z-score"),
    (RGBColor(0x8e,0x44,0xad), "Interactions",   "Vib × Temp\nDrift × Faults"),
    (RGBColor(0x16,0xa0,0x85), "Maint. Urgency", "Faults ×\ncycles since service"),
]
for i, (color, title, desc) in enumerate(fe_blocks):
    row, col = divmod(i, 3)
    lx = 0.3 + col * 4.3
    ty = 1.65 + row * 2.3
    rect(slide, lx, ty, 4.0, 2.05, fill=color)
    txt(slide, title, lx + 0.15, ty + 0.12, 3.7, 0.55,
        size=13.5, bold=True, color=WHITE)
    txt(slide, desc, lx + 0.15, ty + 0.7, 3.7, 1.1, size=11.5, color=WHITE)

rect(slide, 0.3, 6.25, 12.7, 0.65, fill=DARK_BLUE)
txt(slide, "★  fault_roll5_sum = #1 predictor   |   61% of XGBoost's predictive gain   "
    "|   Validated by SHAP analysis",
    0.4, 6.3, 12.5, 0.55, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
footer_bar(slide, 6)

# ─────────────────────────────────────────────────────────────────────────────
# 7. MODEL TRAINING
# ─────────────────────────────────────────────────────────────────────────────
slide = blank(prs)
header(slide, "Training Four Models",
       "Cost-sensitive learning + stratified cross-validation for the 42:1 imbalance")

img(slide, f"{FIGURES}/07_roc_pr_curves.png", 0.2, 1.5, w=8.5, h=4.0)

rect(slide, 9.0, 1.5, 4.1, 4.0, fill=LIGHT_BLUE)
txt(slide, "Model Comparison", 9.15, 1.6, 3.8, 0.42, size=12, bold=True, color=DARK_BLUE)

table_rows = [
    ("Logistic Reg.",  "0.557", "0.857"),
    ("Random Forest",  "0.870", "0.964"),
    ("LightGBM",       "0.936", "0.964"),
    (f"{best_name} ★", "0.963", "1.000"),
]
txt(slide, "Model", 9.15, 2.12, 2.0, 0.38, size=9, bold=True, color=DARK_BLUE)
txt(slide, "PR-AUC", 11.2, 2.12, 0.9, 0.38, size=9, bold=True, color=DARK_BLUE)
txt(slide, "Recall", 12.15, 2.12, 0.85, 0.38, size=9, bold=True, color=DARK_BLUE)
for i, (model, ap, rec) in enumerate(table_rows):
    c = GREEN if i == 3 else DARK_GRAY
    b = (i == 3)
    txt(slide, model,  9.15, 2.6 + i*0.62, 2.0, 0.55, size=9.5, bold=b, color=c)
    txt(slide, ap,     11.2, 2.6 + i*0.62, 0.9, 0.55, size=9.5, bold=b, color=c, align=PP_ALIGN.CENTER)
    txt(slide, rec,    12.15, 2.6 + i*0.62, 0.85, 0.55, size=9.5, bold=b, color=c, align=PP_ALIGN.CENTER)

txt(slide, "Why XGBoost?", 9.15, 5.0, 3.8, 0.38, size=11, bold=True, color=DARK_BLUE)
txt(slide, "Gradient boosting focuses\niterations on hard-to-classify\nfailure events.\n"
    "scale_pos_weight=42 bakes\nthe cost asymmetry into\nthe loss function.",
    9.15, 5.4, 3.8, 1.4, size=10, color=DARK_GRAY)

footer_bar(slide, 7)

# ─────────────────────────────────────────────────────────────────────────────
# 8. RESULTS — ZERO MISSED FAILURES
# ─────────────────────────────────────────────────────────────────────────────
slide = blank(prs)
header(slide, "Results: Zero Missed Failures",
       f"{best_name} on held-out 20% test set (1,200 records | 28 failure events)")

kpi_data = [
    ("0.9992", "ROC-AUC",  DARK_BLUE),
    ("0.9632", "PR-AUC",   MED_BLUE),
    ("0.9180", "F1 Score", GREEN),
    ("1.0000", "Recall",   RED),
]
for i, (val, lbl, bg) in enumerate(kpi_data):
    kpi_box(slide, 0.25 + i*3.1, 1.6, val, lbl, bg=bg)

img(slide, f"{FIGURES}/08_confusion_matrices.png", 0.2, 3.25, w=8.5, h=3.9)

rect(slide, 9.0, 3.25, 4.1, 3.9, fill=LIGHT_BLUE)
txt(slide, "What This Means", 9.15, 3.4, 3.8, 0.45, size=12, bold=True, color=DARK_BLUE)
meaning = [
    (0, "Recall = 1.00: every\n  failure event caught"),
    (0, "Precision = 0.85: ~1 false\n  alarm per 6.5 alerts"),
    (0, "ROC-AUC = 0.999: near-\n  perfect discrimination"),
    (0, "CV variance ±0.023:\n  stable, generalisable"),
]
bullets(slide, meaning, 9.1, 3.95, 3.9, 3.1, size=11, gap=0.77)

footer_bar(slide, 8)

# ─────────────────────────────────────────────────────────────────────────────
# 9. EXPLAINABILITY
# ─────────────────────────────────────────────────────────────────────────────
slide = blank(prs)
header(slide, "Why Does the Model Predict Failure?",
       "SHAP values + decision rules — every alert is explainable")

img(slide, f"{FIGURES}/11_shap_summary.png", 0.15, 1.5, w=8.6, h=4.5)

rect(slide, 8.95, 1.5, 4.15, 4.5, fill=LIGHT_BLUE)
txt(slide, "What Engineers See", 9.1, 1.6, 3.9, 0.45,
    size=12, bold=True, color=DARK_BLUE)
shap_insights = [
    (0, "fault_roll5_sum (61%):\n  rapid fault accumulation"),
    (0, "fault_roll5_mean (21%):\n  sustained fault intensity"),
    (0, "vibration trend: mechanical\n  wear signal"),
    (0, "temperature trend: thermal\n  stress history"),
    (0, "sensor_drift: instrument\n  degradation signal"),
]
bullets(slide, shap_insights, 9.05, 2.15, 4.0, 3.7, size=11, gap=0.72)

rect(slide, 0.15, 6.1, 12.85, 0.65, fill=DARK_BLUE)
txt(slide, "Decision rule for engineers: IF fault_roll5_sum > 9.5 "
    "AND temp_sensor_2_5cycle_avg > 93.8°C  →  FLAG FOR INSPECTION",
    0.25, 6.15, 12.65, 0.55, size=11, bold=True, color=WHITE,
    align=PP_ALIGN.CENTER, font="Courier New")
footer_bar(slide, 9)

# ─────────────────────────────────────────────────────────────────────────────
# 10. THRESHOLD
# ─────────────────────────────────────────────────────────────────────────────
slide = blank(prs)
header(slide, "Choosing the Right Alert Threshold",
       "The model outputs a probability — the threshold is an operational decision")

img(slide, f"{FIGURES}/12_threshold_optimization.png", 0.15, 1.5, w=7.9, h=4.5)

rect(slide, 8.25, 1.5, 4.85, 4.5, fill=LIGHT_BLUE)
txt(slide, "Threshold Guide", 8.4, 1.6, 4.6, 0.45,
    size=12, bold=True, color=DARK_BLUE)

tiers = [
    (RED,    "Safety-critical",    "0.20–0.30",
     "Max recall\nMore inspections"),
    (ORANGE, "Balanced (default)", "0.50–0.68",
     "Optimal F1\nBalanced cost"),
    (GREEN,  "Cost-optimised",     "0.70–0.80",
     "Fewer alerts\nMarginal risk"),
]
for i, (color, ctx, t, desc) in enumerate(tiers):
    rect(slide, 8.35, 2.15 + i*1.22, 4.65, 1.1, fill=color)
    txt(slide, ctx, 8.5, 2.17 + i*1.22, 4.4, 0.38, size=11, bold=True, color=WHITE)
    txt(slide, f"Threshold: {t}", 8.5, 2.53 + i*1.22, 2.5, 0.35, size=10, color=WHITE, bold=True)
    txt(slide, desc, 11.0, 2.17 + i*1.22, 2.0, 0.7, size=9.5, color=WHITE)

txt(slide, "At threshold 0.68: F1 = 0.949 | Recall = 1.00 | Precision = 0.903",
    0.15, 6.1, 12.85, 0.45, size=12, bold=True, color=GREEN,
    align=PP_ALIGN.CENTER)
footer_bar(slide, 10)

# ─────────────────────────────────────────────────────────────────────────────
# 11. DECISION SUPPORT SYSTEM
# ─────────────────────────────────────────────────────────────────────────────
slide = blank(prs)
header(slide, "From Predictions to Action",
       "Risk scores mapped to a four-tier maintenance action framework")

img(slide, f"{FIGURES}/13_risk_scoring.png", 0.15, 1.5, w=8.0, h=3.9)

rect(slide, 8.35, 1.5, 4.75, 3.9, fill=LIGHT_BLUE)
txt(slide, "Risk Tiers", 8.5, 1.6, 4.5, 0.42, size=12, bold=True, color=DARK_BLUE)
tiers_right = [
    (RED,    "🔴 CRITICAL (≥0.60)", "Ground — urgent inspection"),
    (ORANGE, "🟠 HIGH (0.35–0.59)", "Next maintenance window"),
    (GOLD,   "🟡 MEDIUM (0.15–0.34)","Increase monitoring"),
    (GREEN,  "🟢 LOW (<0.15)",       "Standard operations"),
]
for i, (color, tier, action) in enumerate(tiers_right):
    rect(slide, 8.45, 2.1 + i*0.8, 4.55, 0.72, fill=color)
    txt(slide, tier,   8.58, 2.12 + i*0.8, 4.2, 0.33, size=10, bold=True, color=WHITE)
    txt(slide, action, 8.58, 2.43 + i*0.8, 4.2, 0.33, size=9.5, color=WHITE)

# Pipeline flow
rect(slide, 0.15, 5.55, 12.85, 1.25, fill=DARK_BLUE)
txt(slide, "Integration Pipeline", 0.3, 5.6, 3.0, 0.4,
    size=11, bold=True, color=WHITE)
flow = ["1. Ingest\nsensor logs", "2. Compute\nfeatures",
        "3. Score with\nXGBoost", "4. Assign\nrisk tier",
        "5. Dispatch\nwork order"]
for i, step in enumerate(flow):
    txt(slide, step, 0.3 + i*2.55, 5.98, 2.4, 0.72,
        size=10, color=WHITE, align=PP_ALIGN.CENTER)
    if i < 4:
        txt(slide, "→", 2.55 + i*2.55, 6.1, 0.35, 0.45,
            size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

footer_bar(slide, 11)

# ─────────────────────────────────────────────────────────────────────────────
# 12. KEY TAKEAWAYS
# ─────────────────────────────────────────────────────────────────────────────
slide = blank(prs)
rect(slide, 0, 0, 13.33, 7.5, fill=DARK_BLUE)
rect(slide, 0, 5.9, 13.33, 1.6, fill=MED_BLUE)
rect(slide, 3.5, 1.05, 6.3, 0.07, fill=WHITE)

txt(slide, "Key Takeaways",
    0.5, 0.22, 12.3, 0.75, size=38, bold=True, color=WHITE,
    align=PP_ALIGN.CENTER, font="Calibri Light")

kpi_final = [
    ("0.9992", "ROC-AUC",  DARK_BLUE),
    ("0.9632", "PR-AUC",   MED_BLUE),
    ("1.0000", "Recall",   RED),
    ("0.9180", "F1 Score", GREEN),
]
for i, (val, lbl, bg) in enumerate(kpi_final):
    rect(slide, 0.4 + i*3.1, 1.22, 2.85, 1.45, fill=bg)
    txt(slide, val, 0.4 + i*3.1, 1.27, 2.85, 0.85,
        size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(slide, lbl, 0.4 + i*3.1, 2.1, 2.85, 0.45,
        size=11, color=WHITE, align=PP_ALIGN.CENTER)

takeaways = [
    "XGBoost achieves perfect recall — zero failures missed on the test set",
    "Rolling fault accumulation (fault_roll5_sum) is the decisive signal — 61% of model gain",
    "A shallow decision tree provides human-readable rules for maintenance engineers",
    "The inference pipeline (predict.py) is production-ready for integration with FDMS",
]
for i, t in enumerate(takeaways):
    txt(slide, f"✓   {t}", 0.4, 2.92 + i * 0.62, 12.5, 0.58, size=13.5, color=WHITE)

txt(slide, "Prepared for RCAA  |  Data Science Assessment  |  March 2026",
    0.5, 6.08, 12.3, 0.45, size=12, color=WHITE, align=PP_ALIGN.CENTER)
txt(slide, "Questions welcome — I can walk through any section in detail",
    0.5, 6.6, 12.3, 0.45, size=11, color=LIGHT_BLUE,
    align=PP_ALIGN.CENTER, italic=True)

# ── Save ──────────────────────────────────────────────────────────────────────
output = "RCAA_Predictive_Maintenance_Slides.pptx"
prs.save(output)
print(f"Slides saved: {output}")
